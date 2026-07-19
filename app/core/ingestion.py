import json
import re
import logging
from pathlib import Path
from typing import List, Dict, Any
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from bs4 import BeautifulSoup
import pandas as pd
from app.config import settings

logger = logging.getLogger("rag_system.ingestion")

def clean_text(text: str) -> str:
    """Removes redundant whitespace and normalizes text content."""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def detect_section_header(text: str, current_section: str) -> str:
    """Attempts to detect section headers in a text line/segment."""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    for line in lines[:3]:  # Check first few lines of a chunk
        # Regex for numbered headers: e.g., "1. Introduction" or "1.2.3 Data Collection"
        if re.match(r'^([I|V|X|L|C]+|\d+(\.\d+)*)\.?\s+[A-Z][a-zA-Z\s]{2,50}$', line):
            return line
        # Regex for typical uppercase headers: e.g., "ABSTRACT", "INTRODUCTION"
        if re.match(r'^[A-Z\s]{4,30}$', line):
            return line
    return current_section

class DocumentIngestor:
    @staticmethod
    def load_document(file_path: Path) -> List[Dict[str, Any]]:
        """
        Loads document content and extracts text and metadata.
        Returns a list of dicts: [{'text': str, 'metadata': {'source': str, 'page': int, 'section': str}}]
        """
        file_path = Path(file_path)
        ext = file_path.suffix.lower()
        logger.info(f"Ingesting document: {file_path.name} (type: {ext})")
        
        try:
            if ext == '.pdf':
                return DocumentIngestor._load_pdf(file_path)
            elif ext == '.docx':
                return DocumentIngestor._load_docx(file_path)
            elif ext in ['.pptx', '.ppt']:
                return DocumentIngestor._load_pptx(file_path)
            elif ext in ['.html', '.htm']:
                return DocumentIngestor._load_html(file_path)
            elif ext == '.csv':
                return DocumentIngestor._load_csv(file_path)
            elif ext == '.json':
                return DocumentIngestor._load_json(file_path)
            elif ext in ['.txt', '.md']:
                return DocumentIngestor._load_txt_md(file_path)
            else:
                logger.error(f"Ingestion failed: Unsupported extension '{ext}' for file {file_path.name}")
                raise ValueError(f"Unsupported file format: {ext}")
        except Exception as e:
            logger.error(f"Error parsing file {file_path.name}: {e}", exc_info=True)
            raise

    @staticmethod
    def _load_pdf(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        doc = fitz.open(str(file_path))
        current_section = "Abstract"
        
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            cleaned = clean_text(text)
            if not cleaned:
                continue
            
            current_section = detect_section_header(text, current_section)
            results.append({
                "text": cleaned,
                "metadata": {
                    "source": file_path.name,
                    "page": page_num,
                    "section": current_section
                }
            })
        doc.close()
        return results

    @staticmethod
    def _load_docx(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        doc = docx.Document(str(file_path))
        current_section = "Introduction"
        
        current_block = []
        current_block_len = 0
        page_num = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            if para.style.name.startswith('Heading') or re.match(r'^\d+(\.\d+)*\s+[A-Z]', text):
                current_section = text
            
            current_block.append(text)
            current_block_len += len(text)
            
            if current_block_len >= 2500:
                block_text = " ".join(current_block)
                results.append({
                    "text": clean_text(block_text),
                    "metadata": {
                        "source": file_path.name,
                        "page": page_num,
                        "section": current_section
                    }
                })
                current_block = []
                current_block_len = 0
                page_num += 1
                
        if current_block:
            block_text = " ".join(current_block)
            results.append({
                "text": clean_text(block_text),
                "metadata": {
                    "source": file_path.name,
                    "page": page_num,
                    "section": current_section
                }
            })
            
        return results

    @staticmethod
    def _load_pptx(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        prs = Presentation(str(file_path))
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            title = f"Slide {slide_num}"
            
            if slide.shapes.title:
                title = clean_text(slide.shapes.title.text)
                
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            full_text = "\n".join(slide_texts)
            cleaned = clean_text(full_text)
            if not cleaned:
                continue
                
            results.append({
                "text": cleaned,
                "metadata": {
                    "source": file_path.name,
                    "page": slide_num,
                    "section": title
                }
            })
        return results

    @staticmethod
    def _load_html(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            
        for script in soup(["script", "style"]):
            script.decompose()
            
        current_section = "Web Page Content"
        
        body = soup.body if soup.body else soup
        blocks = []
        current_block = []
        current_len = 0
        
        for element in body.descendants:
            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                if current_block:
                    blocks.append((clean_text(" ".join(current_block)), current_section))
                    current_block = []
                    current_len = 0
                current_section = element.get_text().strip()
            elif element.name in ['p', 'div', 'li'] and element.string:
                txt = element.string.strip()
                if txt:
                    current_block.append(txt)
                    current_len += len(txt)
                    if current_len >= 2500:
                        blocks.append((clean_text(" ".join(current_block)), current_section))
                        current_block = []
                        current_len = 0
                        
        if current_block:
            blocks.append((clean_text(" ".join(current_block)), current_section))
            
        for idx, (text, section) in enumerate(blocks, 1):
            if text:
                results.append({
                    "text": text,
                    "metadata": {
                        "source": file_path.name,
                        "page": idx,
                        "section": section
                    }
                })
        return results

    @staticmethod
    def _load_csv(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        df = pd.read_csv(file_path)
        
        row_strings = []
        for idx, row in df.iterrows():
            items = [f"{col}: {val}" for col, val in row.items() if pd.notna(val)]
            row_strings.append(f"Row {idx+1}: " + ", ".join(items))
            
        chunk_size = 10
        for i in range(0, len(row_strings), chunk_size):
            chunk_rows = row_strings[i:i+chunk_size]
            text = "\n".join(chunk_rows)
            results.append({
                "text": text,
                "metadata": {
                    "source": file_path.name,
                    "page": (i // chunk_size) + 1,
                    "section": "Data Table"
                }
            })
        return results

    @staticmethod
    def _load_json(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            
        if isinstance(data, list):
            for idx, item in enumerate(data, 1):
                item_str = json.dumps(item, indent=2)
                results.append({
                    "text": f"Item {idx}:\n{item_str}",
                    "metadata": {
                        "source": file_path.name,
                        "page": idx,
                        "section": "Structured List"
                    }
                })
        elif isinstance(data, dict):
            for idx, (key, val) in enumerate(data.items(), 1):
                val_str = json.dumps(val, indent=2)
                results.append({
                    "text": f"Key: {key}\nValue:\n{val_str}",
                    "metadata": {
                        "source": file_path.name,
                        "page": idx,
                        "section": f"Section: {key}"
                    }
                })
        return results

    @staticmethod
    def _load_txt_md(file_path: Path) -> List[Dict[str, Any]]:
        results = []
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
            
        sections = re.split(r'^(#+\s+.*)$', content, flags=re.MULTILINE)
        current_section = "General"
        page_num = 1
        
        if len(sections) <= 1:
            lines = content.split('\n')
            current_block = []
            current_len = 0
            for line in lines:
                current_block.append(line)
                current_len += len(line)
                if current_len >= 2500:
                    text = "\n".join(current_block)
                    results.append({
                        "text": clean_text(text),
                        "metadata": {
                            "source": file_path.name,
                            "page": page_num,
                            "section": current_section
                        }
                    })
                    current_block = []
                    current_len = 0
                    page_num += 1
            if current_block:
                text = "\n".join(current_block)
                results.append({
                    "text": clean_text(text),
                    "metadata": {
                        "source": file_path.name,
                        "page": page_num,
                        "section": current_section
                    }
                })
        else:
            current_text = ""
            for item in sections:
                if item.startswith('#'):
                    if current_text.strip():
                        results.append({
                            "text": clean_text(current_text),
                            "metadata": {
                                "source": file_path.name,
                                "page": page_num,
                                "section": current_section
                            }
                        })
                        page_num += 1
                    current_section = item.strip().lstrip('#').strip()
                    current_text = ""
                else:
                    current_text += item
            if current_text.strip():
                results.append({
                    "text": clean_text(current_text),
                    "metadata": {
                        "source": file_path.name,
                        "page": page_num,
                        "section": current_section
                    }
                })
        return results

class RecursiveMetadataChunker:
    def __init__(self, chunk_size: int = None, chunk_overlap: int = None):
        self.chunk_size = chunk_size or settings.CHUNK_SIZE
        self.chunk_overlap = chunk_overlap or settings.CHUNK_OVERLAP

    def chunk_documents(self, documents: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Splits documents into smaller overlapping chunks while retaining parent metadata.
        """
        chunks = []
        for doc in documents:
            text = doc["text"]
            metadata = doc["metadata"]
            
            if len(text) <= self.chunk_size:
                chunks.append({
                    "text": text,
                    "metadata": metadata.copy()
                })
                continue
                
            start = 0
            while start < len(text):
                end = start + self.chunk_size
                if end < len(text):
                    space_idx = text.rfind(' ', end - 30, end)
                    if space_idx != -1:
                        end = space_idx
                
                chunk_text = text[start:end].strip()
                if chunk_text:
                    chunks.append({
                        "text": chunk_text,
                        "metadata": metadata.copy()
                    })
                    
                start = end - self.chunk_overlap
                if start >= len(text) - self.chunk_overlap:
                    break
        logger.info(f"Chunked document into {len(chunks)} fragments (size={self.chunk_size}, overlap={self.chunk_overlap})")
        return chunks
